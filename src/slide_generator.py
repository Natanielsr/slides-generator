from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from src.slide_data import SlideData
from src.font_size_calculation import FontSizeCalculation

class SlideGenerator:
    def __init__(self, presentation, slide_layout, slide_data : SlideData):
        self.__presentation = presentation
        self.__slide_layout = slide_layout
        self.__slide_data = slide_data

        self.calculate_font_size()

    def create_stanza_slide(self, insert_title : bool = False):
        self.__slide_pptx = self.__presentation.slides.add_slide(self.__slide_layout)
        self.set_background_color()
        if(insert_title):
            self.add_title()
        self.add_content()

    def set_background_color(self):
        # Define o fundo preto do slide
        background = self.__slide_pptx.background
        fill = background.fill
        fill.solid()  # Configura para cor sólida
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Preto

    def add_title(self):
        title = self.__slide_pptx.shapes.title
        title.text = self.__slide_data.title

    def add_content(self):

        left = top = Inches(1)  # Ajuste as coordenadas conforme necessário
        width = Inches(8)  # Largura fixa do textbox
        height = Inches(2)  # Altura fixa do textbox

        txBox = self.__slide_pptx.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = self.__slide_data.stanza

        for p in tf.paragraphs:
            # Define a fonte e a cor (como no seu código anterior)
            p.font.size = self.__stanza_font_size
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Times New Roman'

         # Habilita a quebra de linha automática
        tf.word_wrap = True
   
    def calculate_font_size(self):
        fsc = FontSizeCalculation()
        sfs = fsc.calcular_tamanho_fonte(self.__slide_data.stanza)
        self.__stanza_font_size = Pt(sfs)