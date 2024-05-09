from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class SlidesGenerator:

    def __init__(self, presentation, layout, font_size = 36):     
        self.__FONT_SIZE = font_size
        self.__title_shape = None
        self.__content_shape = None
        self.__presentation = presentation
        self.__layout = layout
        
    def create_stanza_slide(self, title, stanza):
        slide = self._add_slide_with_layout()

        self._set_slide_title(slide, title)
        self._format_slide_title()

        self._set_slide_content(slide, stanza)
        self._format_slide_content()
    
    
    def _add_slide_with_layout(self):
        return self.__presentation.slides.add_slide(self.__layout)
    
    def _set_slide_title(self, slide, title):
        # Adicionando texto ao título do slide
        title_shape = slide.shapes.title
        title_shape.text = title

        self.__title_shape = title_shape
    
    def _format_slide_title(self):
        self.__title_shape.text_frame.paragraphs[0].font.bold = True


    def _set_slide_content(self, slide, stanza_text):
        # Adicionando texto ao corpo do slide
        content_shape = slide.placeholders[1]  # O índice 1 representa o espaço reservado para o corpo do slide
        content_shape.text = stanza_text

        self.__content_shape = content_shape
        
    def _format_slide_content(self):
        text_frame = self.__content_shape.text_frame
        for p in text_frame.paragraphs:
            p.font.size = Pt(self.__FONT_SIZE)
            p.level = 0  # Define o nível de lista como 0
            p.alignment = PP_ALIGN.CENTER  # Alinha o texto ao centro
            p.font.color.rgb = RGBColor(0, 0, 0)  # Define a cor do texto como preto